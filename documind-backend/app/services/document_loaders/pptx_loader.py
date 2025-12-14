"""
PowerPoint (PPTX) document loader
"""

from typing import Dict, Any, List
from pptx import Presentation
from pptx.shapes.base import BaseShape
import structlog

from .base import DocumentLoader, DocumentContent, LoaderError

logger = structlog.get_logger(__name__)


class PPTXLoader(DocumentLoader):
    """Loader for PowerPoint (PPTX) documents"""
    
    def load(self) -> DocumentContent:
        """
        Load and extract content from PPTX file
        
        Returns:
            DocumentContent: Extracted content with text, slides, tables, and metadata
            
        Raises:
            LoaderError: If PPTX loading fails
        """
        try:
            logger.info("loading_pptx", file_path=str(self.file_path))
            
            metadata = self._extract_metadata()
            slides_data = []
            tables_data = []
            full_text = []
            
            # Open presentation
            prs = Presentation(self.file_path)
            
            # Extract presentation properties
            core_props = prs.core_properties
            metadata.update({
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "comments": core_props.comments or "",
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "last_modified_by": core_props.last_modified_by or "",
                "revision": core_props.revision or 0,
                "slide_count": len(prs.slides),
            })
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text_parts = []
                slide_tables = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    shape_text = self._extract_shape_text(shape)
                    if shape_text:
                        slide_text_parts.append(shape_text)
                    
                    # Extract tables
                    if hasattr(shape, "table"):
                        table_data = self._extract_table(shape.table)
                        if table_data:
                            slide_tables.append(table_data)
                            tables_data.append(table_data)
                
                slide_text = "\n".join(slide_text_parts)
                normalized_slide_text = self._normalize_text(slide_text)
                
                slides_data.append({
                    "slide_number": slide_num,
                    "text": normalized_slide_text,
                    "char_count": len(normalized_slide_text),
                    "table_count": len(slide_tables),
                })
                
                if normalized_slide_text:
                    full_text.append(f"Slide {slide_num}:\n{normalized_slide_text}")
            
            # Combine all text
            combined_text = "\n\n".join(full_text)
            
            # Apply preprocessing
            from .preprocessing import preprocess_text
            preprocess_result = preprocess_text(
                combined_text,
                remove_page_nums=True,
                remove_headers_footers=True,
                detect_lang=True
            )
            combined_text = preprocess_result["text"]
            combined_text = self._normalize_text(combined_text, apply_preprocessing=False)
            
            # Add preprocessing metadata
            if "page_numbers_removed" in preprocess_result:
                metadata["page_numbers_removed"] = preprocess_result["page_numbers_removed"]
            if "headers_footers_removed" in preprocess_result:
                metadata["headers_footers_removed"] = preprocess_result["headers_footers_removed"]
            if "language_detection" in preprocess_result:
                lang_info = preprocess_result["language_detection"]
                metadata["detected_language"] = lang_info.get("language", "unknown")
                metadata["language_confidence"] = lang_info.get("confidence", 0.0)
            
            metadata.update({
                "table_count": len(tables_data),
                "total_char_count": len(combined_text),
                "extraction_method": "python-pptx",
            })
            
            logger.info(
                "pptx_loaded_successfully",
                file_path=str(self.file_path),
                slide_count=len(prs.slides),
                table_count=len(tables_data),
                char_count=len(combined_text)
            )
            
            return DocumentContent(
                text=combined_text,
                metadata=metadata,
                pages=slides_data,  # Use pages field for slides
                tables=tables_data
            )
            
        except Exception as e:
            error_msg = f"Failed to load PPTX file: {str(e)}"
            logger.error("pptx_load_error", file_path=str(self.file_path), error=str(e))
            raise LoaderError(error_msg, file_path=str(self.file_path), original_error=e)
    
    def _extract_shape_text(self, shape: BaseShape) -> str:
        """
        Extract text from a shape
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            str: Extracted text
        """
        try:
            if hasattr(shape, "text"):
                return shape.text
            elif hasattr(shape, "text_frame"):
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)
                return "\n".join(text_parts)
            return ""
        except Exception:
            return ""
    
    def _extract_table(self, table) -> Dict[str, Any]:
        """
        Extract data from a table
        
        Args:
            table: python-pptx Table object
            
        Returns:
            Dict[str, Any]: Table data
        """
        try:
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() if cell.text else "" for cell in row.cells]
                rows_data.append(row_data)
            
            return {
                "rows": rows_data,
                "row_count": len(rows_data),
                "column_count": len(rows_data[0]) if rows_data else 0,
            }
        except Exception:
            return {}
    
    def supports_streaming(self) -> bool:
        """PPTX loading supports streaming for large files"""
        return True
    
    def load_streaming(self, max_slides: int = None) -> DocumentContent:
        """
        Load PPTX file with streaming support
        
        Args:
            max_slides: Maximum number of slides to load (None for all)
            
        Returns:
            DocumentContent: Extracted content
        """
        try:
            logger.info("loading_pptx_streaming", file_path=str(self.file_path), max_slides=max_slides)
            
            metadata = self._extract_metadata()
            slides_data = []
            tables_data = []
            full_text = []
            
            prs = Presentation(self.file_path)
            total_slides = len(prs.slides)
            slides_to_load = min(max_slides or total_slides, total_slides)
            
            metadata.update({
                "slide_count": total_slides,
                "slides_loaded": slides_to_load,
                "is_partial": slides_to_load < total_slides,
            })
            
            for slide_num in range(1, slides_to_load + 1):
                slide = prs.slides[slide_num - 1]
                slide_text_parts = []
                
                for shape in slide.shapes:
                    shape_text = self._extract_shape_text(shape)
                    if shape_text:
                        slide_text_parts.append(shape_text)
                    
                    if hasattr(shape, "table"):
                        table_data = self._extract_table(shape.table)
                        if table_data:
                            tables_data.append(table_data)
                
                slide_text = "\n".join(slide_text_parts)
                normalized_slide_text = self._normalize_text(slide_text)
                
                slides_data.append({
                    "slide_number": slide_num,
                    "text": normalized_slide_text,
                    "char_count": len(normalized_slide_text),
                })
                
                if normalized_slide_text:
                    full_text.append(f"Slide {slide_num}:\n{normalized_slide_text}")
            
            combined_text = "\n\n".join(full_text)
            
            # Apply preprocessing
            from .preprocessing import preprocess_text
            preprocess_result = preprocess_text(
                combined_text,
                remove_page_nums=True,
                remove_headers_footers=True,
                detect_lang=True
            )
            combined_text = preprocess_result["text"]
            combined_text = self._normalize_text(combined_text, apply_preprocessing=False)
            
            # Add preprocessing metadata
            if "page_numbers_removed" in preprocess_result:
                metadata["page_numbers_removed"] = preprocess_result["page_numbers_removed"]
            if "headers_footers_removed" in preprocess_result:
                metadata["headers_footers_removed"] = preprocess_result["headers_footers_removed"]
            if "language_detection" in preprocess_result:
                lang_info = preprocess_result["language_detection"]
                metadata["detected_language"] = lang_info.get("language", "unknown")
                metadata["language_confidence"] = lang_info.get("confidence", 0.0)
            
            metadata.update({
                "table_count": len(tables_data),
                "total_char_count": len(combined_text),
                "extraction_method": "python-pptx_streaming",
            })
            
            return DocumentContent(
                text=combined_text,
                metadata=metadata,
                pages=slides_data,
                tables=tables_data
            )
            
        except Exception as e:
            error_msg = f"Failed to stream PPTX file: {str(e)}"
            logger.error("pptx_streaming_error", file_path=str(self.file_path), error=str(e))
            raise LoaderError(error_msg, file_path=str(self.file_path), original_error=e)

